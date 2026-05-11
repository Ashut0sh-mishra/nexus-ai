"""Phase 6S — PPTX import endpoint tests.

Fixture decks are produced at test time with python-pptx so the suite
remains binary-clean (no checked-in .pptx files). Every test uses an
in-memory SQLite database via the same dependency-override pattern as
``test_lifecycle_route.py``.
"""

from __future__ import annotations

import io
from typing import Any, AsyncIterator

import httpx
import pytest
from pptx import Presentation
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession, async_sessionmaker, create_async_engine

from agent.slide_schema import validate_deck
from database.connection import get_db
from database.models import Base, SlideDeck, Task
from main import app


def _make_pptx_bytes(slides: list[dict[str, Any]]) -> bytes:
    prs = Presentation()
    layout = prs.slide_layouts[1]  # title + content
    for spec in slides:
        s = prs.slides.add_slide(layout)
        if s.shapes.title is not None:
            s.shapes.title.text = spec.get("title", "") or ""
        body_ph = None
        for ph in s.placeholders:
            try:
                if ph.placeholder_format.idx == 1:
                    body_ph = ph
                    break
            except Exception:
                pass
        if body_ph is not None and spec.get("body"):
            tf = body_ph.text_frame
            tf.text = spec["body"][0]
            for line in spec["body"][1:]:
                p = tf.add_paragraph()
                p.text = line
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


async def _setup_app():
    engine = create_async_engine("sqlite+aiosqlite:///:memory:", future=True)
    async with engine.begin() as conn:
        await conn.run_sync(Base.metadata.create_all)
    Session = async_sessionmaker(engine, class_=AsyncSession, expire_on_commit=False)

    async def _override() -> AsyncIterator[AsyncSession]:
        async with Session() as s:
            yield s

    app.dependency_overrides[get_db] = _override
    return Session, engine


async def _teardown(engine):
    app.dependency_overrides.clear()
    await engine.dispose()


def _client() -> httpx.AsyncClient:
    return httpx.AsyncClient(
        transport=httpx.ASGITransport(app=app), base_url="http://t"
    )


@pytest.mark.asyncio
async def test_import_fixture_pptx_succeeds_and_persists():
    Session, engine = await _setup_app()
    try:
        pptx = _make_pptx_bytes(
            [
                {"title": "Quarterly Review", "body": ["Theme: scale with discipline"]},
                {"title": "Highlights", "body": ["Revenue up 38%", "ARR crossed $42M"]},
                {"title": "Next steps", "body": ["Hire two AEs", "Launch SOC2 audit"]},
            ]
        )
        async with _client() as c:
            r = await c.post(
                "/api/import/pptx",
                files={
                    "file": (
                        "deck.pptx",
                        pptx,
                        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    )
                },
            )
        assert r.status_code == 200, r.text
        body = r.json()
        assert body["slide_count"] == 3
        assert len(body["slides"]) == 3
        assert body["slides"][0]["layout"] == "title"
        assert body["slides"][0]["title"] == "Quarterly Review"
        assert body["source"]["source_slide_count"] == 3
        assert body["source"]["truncated"] is False

        async with Session() as s:
            task = (
                await s.execute(select(Task).where(Task.id == body["task_id"]))
            ).scalar_one()
            deck = (
                await s.execute(
                    select(SlideDeck).where(SlideDeck.task_id == body["task_id"])
                )
            ).scalar_one()
            assert task.status == "done"
            assert task.progress_pct == 100.0
            assert task.topic == "Quarterly Review"
            assert deck.slide_count == 3
    finally:
        await _teardown(engine)


@pytest.mark.asyncio
async def test_import_preserves_slide_count_and_titles():
    _, engine = await _setup_app()
    try:
        pptx = _make_pptx_bytes(
            [
                {"title": "Slide A", "body": ["a1"]},
                {"title": "Slide B", "body": ["b1", "b2"]},
                {"title": "Slide C", "body": ["c1"]},
                {"title": "Slide D", "body": ["d1"]},
                {"title": "Slide E", "body": ["e1"]},
            ]
        )
        async with _client() as c:
            r = await c.post(
                "/api/import/pptx",
                files={"file": ("five.pptx", pptx, "application/octet-stream")},
            )
        assert r.status_code == 200
        body = r.json()
        assert body["slide_count"] == 5
        assert [s["title"] for s in body["slides"]] == [
            "Slide A",
            "Slide B",
            "Slide C",
            "Slide D",
            "Slide E",
        ]
    finally:
        await _teardown(engine)


@pytest.mark.asyncio
async def test_import_output_validates_against_slide_schema():
    _, engine = await _setup_app()
    try:
        pptx = _make_pptx_bytes(
            [
                {"title": "Hello", "body": ["world"]},
                {"title": "Empty body", "body": []},
                {"title": "", "body": ["only body"]},
            ]
        )
        async with _client() as c:
            r = await c.post(
                "/api/import/pptx",
                files={"file": ("v.pptx", pptx, "application/octet-stream")},
            )
        assert r.status_code == 200
        slides = r.json()["slides"]
        results = validate_deck(slides)
        assert all(res.ok for res in results), [
            [e.to_dict() for e in r.errors] for r in results if not r.ok
        ]
    finally:
        await _teardown(engine)


@pytest.mark.asyncio
async def test_import_rejects_non_pptx_extension():
    _, engine = await _setup_app()
    try:
        async with _client() as c:
            r = await c.post(
                "/api/import/pptx",
                files={"file": ("not-a-deck.txt", b"hello", "text/plain")},
            )
        assert r.status_code == 400
        assert r.json()["detail"]["error"] == "bad_extension"
    finally:
        await _teardown(engine)


@pytest.mark.asyncio
async def test_import_rejects_corrupt_pptx():
    _, engine = await _setup_app()
    try:
        async with _client() as c:
            r = await c.post(
                "/api/import/pptx",
                files={
                    "file": (
                        "fake.pptx",
                        b"this is not a real pptx archive",
                        "application/octet-stream",
                    )
                },
            )
        assert r.status_code == 400
        assert r.json()["detail"]["error"] == "corrupt"
    finally:
        await _teardown(engine)


@pytest.mark.asyncio
async def test_import_rejects_oversize_file(monkeypatch):
    _, engine = await _setup_app()
    try:
        import api.routes.import_pptx as route_mod

        monkeypatch.setattr(route_mod, "MAX_BYTES", 1024, raising=True)
        huge = b"\x00" * (route_mod.MAX_BYTES + 256)
        async with _client() as c:
            r = await c.post(
                "/api/import/pptx",
                files={"file": ("big.pptx", huge, "application/octet-stream")},
            )
        assert r.status_code == 400
        assert r.json()["detail"]["error"] == "too_large"
    finally:
        await _teardown(engine)


@pytest.mark.asyncio
async def test_imported_deck_is_editable_via_existing_slides_get():
    """Phase 6S contract: imports show up via the same /api/slides path."""
    _, engine = await _setup_app()
    try:
        pptx = _make_pptx_bytes([{"title": "Alpha", "body": ["one", "two"]}])
        async with _client() as c:
            r = await c.post(
                "/api/import/pptx",
                files={"file": ("a.pptx", pptx, "application/octet-stream")},
            )
            assert r.status_code == 200
            task_id = r.json()["task_id"]
            g = await c.get(f"/api/slides/{task_id}")
        assert g.status_code == 200
        body = g.json()
        assert body["task_id"] == task_id
        assert body["slide_count"] == 1
        assert body["slides"][0]["layout"] == "title"
        assert body["slides"][0]["title"] == "Alpha"
    finally:
        await _teardown(engine)
