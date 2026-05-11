"""Phase 6L-UX-Fix — server-side persistence of edited decks.

These tests cover the new ``PUT /api/slides/{task_id}`` endpoint and prove
that downstream surfaces (``GET /api/slides/{task_id}``, ``GET /api/share/{token}``,
PPTX export) all read from the same ``SlideDeck`` row that PUT updates.

The tests use ``httpx.AsyncClient`` with ``ASGITransport`` so the FastAPI
lifespan does not run (no provider keys required), and override
``get_db`` with an in-memory SQLite engine so they never touch a real DB.
"""

from __future__ import annotations

import asyncio
import io
from typing import AsyncIterator

import httpx
import pytest
from pptx import Presentation
from sqlalchemy import select
from sqlalchemy.ext.asyncio import AsyncSession, async_sessionmaker, create_async_engine

from database.connection import get_db
from database.models import Base, ShareToken, SlideDeck, Task
from main import app


# ── canonical fixtures ────────────────────────────────────────────────────


def _seed_deck() -> list[dict]:
    """A small, contract-valid 3-slide deck."""
    return [
        {
            "id": "slide-000",
            "layout": "title",
            "title": "Original Title",
            "subtitle": "Original Subtitle",
            "eyebrow": "Presentation",
        },
        {
            "id": "slide-001",
            "layout": "bullets",
            "title": "Original Bullets",
            "bullets": ["one", "two", "three"],
        },
        {
            "id": "slide-002",
            "layout": "closing",
            "title": "Original Closing",
            "subtitle": "",
            "cta": "Q&A",
        },
    ]


def _edited_deck() -> list[dict]:
    """Same shape as ``_seed_deck`` but with user-edited text."""
    return [
        {
            "id": "slide-000",
            "layout": "title",
            "title": "EDITED Title XYZ",
            "subtitle": "EDITED Subtitle",
            "eyebrow": "Presentation",
        },
        {
            "id": "slide-001",
            "layout": "bullets",
            "title": "EDITED Bullets",
            "bullets": ["alpha edited", "beta edited"],
        },
        {
            "id": "slide-002",
            "layout": "closing",
            "title": "EDITED Closing",
            "subtitle": "",
            "cta": "Thanks",
        },
    ]


async def _setup_app() -> tuple[async_sessionmaker, "object", str]:
    engine = create_async_engine("sqlite+aiosqlite:///:memory:", future=True)
    async with engine.begin() as conn:
        await conn.run_sync(Base.metadata.create_all)
    Session = async_sessionmaker(engine, class_=AsyncSession, expire_on_commit=False)

    async def _get_db_override() -> AsyncIterator[AsyncSession]:
        async with Session() as s:
            yield s

    app.dependency_overrides[get_db] = _get_db_override

    # Seed a completed task with the original deck.
    task_id = "task-edit-test"
    async with Session() as s:
        s.add(
            Task(
                id=task_id,
                topic="Renewables",
                slide_count=3,
                theme="Editorial",
                status="done",
            )
        )
        s.add(
            SlideDeck(
                task_id=task_id,
                slide_data=_seed_deck(),
                theme="Editorial",
                slide_count=3,
            )
        )
        await s.commit()

    return Session, engine, task_id


async def _teardown_app(engine) -> None:
    app.dependency_overrides.clear()
    await engine.dispose()


def _client() -> httpx.AsyncClient:
    transport = httpx.ASGITransport(app=app)
    return httpx.AsyncClient(transport=transport, base_url="http://t")


# ── 1. PUT saves edited slides ────────────────────────────────────────────


def test_put_slides_persists_edits():
    async def _go():
        Session, engine, task_id = await _setup_app()
        try:
            async with _client() as c:
                r = await c.put(
                    f"/api/slides/{task_id}",
                    json={"slides": _edited_deck()},
                )
            assert r.status_code == 200, r.text
            body = r.json()
            assert body["task_id"] == task_id
            assert body["slide_count"] == 3
            titles = [s["title"] for s in body["slides"]]
            assert titles == ["EDITED Title XYZ", "EDITED Bullets", "EDITED Closing"]

            # Row was updated in place.
            async with Session() as s:
                deck = (
                    await s.execute(select(SlideDeck).where(SlideDeck.task_id == task_id))
                ).scalar_one()
                assert deck.slide_count == 3
                assert deck.slide_data[0]["title"] == "EDITED Title XYZ"
        finally:
            await _teardown_app(engine)

    asyncio.run(_go())


# ── 2. GET after PUT returns edited slides ────────────────────────────────


def test_get_after_put_returns_edited_slides():
    async def _go():
        Session, engine, task_id = await _setup_app()
        try:
            async with _client() as c:
                put = await c.put(
                    f"/api/slides/{task_id}",
                    json={"slides": _edited_deck(), "theme": "Vellum"},
                )
                assert put.status_code == 200, put.text
                get = await c.get(f"/api/slides/{task_id}")
            assert get.status_code == 200, get.text
            body = get.json()
            assert body["theme"] == "Vellum"
            assert body["slides"][0]["title"] == "EDITED Title XYZ"
            assert body["slides"][1]["bullets"] == ["alpha edited", "beta edited"]
            # Quality report still attached on read.
            assert "deck_quality" in body
        finally:
            await _teardown_app(engine)

    asyncio.run(_go())


# ── 3. Invalid PUT returns 400 and preserves original deck ────────────────


def test_put_invalid_returns_400_and_does_not_overwrite():
    async def _go():
        Session, engine, task_id = await _setup_app()
        try:
            bad_deck = [
                {
                    "id": "slide-000",
                    "layout": "bullets",
                    # missing required ``title`` and ``bullets`` is empty
                    "bullets": [],
                }
            ]
            async with _client() as c:
                r = await c.put(
                    f"/api/slides/{task_id}",
                    json={"slides": bad_deck},
                )
            assert r.status_code == 400, r.text
            detail = r.json()["detail"]
            assert detail["error"] == "invalid_deck"
            assert isinstance(detail["invalid_slides"], list)
            assert len(detail["invalid_slides"]) == 1
            # Original deck is untouched.
            async with Session() as s:
                deck = (
                    await s.execute(select(SlideDeck).where(SlideDeck.task_id == task_id))
                ).scalar_one()
                titles = [sl["title"] for sl in deck.slide_data]
                assert titles == ["Original Title", "Original Bullets", "Original Closing"]
                assert deck.slide_count == 3
        finally:
            await _teardown_app(engine)

    asyncio.run(_go())


# ── 4. Share link returns edited slides ───────────────────────────────────


def test_share_link_returns_edited_slides():
    async def _go():
        Session, engine, task_id = await _setup_app()
        try:
            async with _client() as c:
                put = await c.put(
                    f"/api/slides/{task_id}",
                    json={"slides": _edited_deck()},
                )
                assert put.status_code == 200, put.text
            # Insert a ShareToken row directly with no expiry so we sidestep
            # a pre-existing SQLite naive-datetime quirk in
            # ``share.view_share``. We are testing whether the share endpoint
            # serves the *edited* deck, not its expiry semantics.
            async with Session() as s:
                s.add(ShareToken(token="tkn-edit", task_id=task_id, expires_at=None))
                await s.commit()
            async with _client() as c:
                view = await c.get("/api/share/tkn-edit")
            assert view.status_code == 200, view.text
            body = view.json()
            assert body["slides"][0]["title"] == "EDITED Title XYZ"
            assert body["slide_count"] == 3
        finally:
            await _teardown_app(engine)

    asyncio.run(_go())


# ── 5. PPTX export uses edited slides ─────────────────────────────────────


def test_pptx_export_includes_edited_text():
    """After PUT, the PPTX renderer reads the edited ``SlideDeck.slide_data``.

    We invoke ``ExportService._export_pptx_sync`` directly with the
    persisted deck to avoid coupling the test to the storage backend or
    the ``/api/export/pptx`` route's StorageService side effects.
    """

    async def _go():
        Session, engine, task_id = await _setup_app()
        try:
            async with _client() as c:
                put = await c.put(
                    f"/api/slides/{task_id}",
                    json={"slides": _edited_deck()},
                )
                assert put.status_code == 200, put.text

            # Read the persisted deck back from the DB and render it.
            async with Session() as s:
                deck = (
                    await s.execute(select(SlideDeck).where(SlideDeck.task_id == task_id))
                ).scalar_one()
                persisted = deck.slide_data

            # Use ExportService with a stub storage so we get bytes back.
            from services.export_service import ExportService

            class _StubStorage:
                def __init__(self) -> None:
                    self.last_data: bytes | None = None

                def put(self, filename: str, data: bytes, content_type: str = "") -> str:
                    self.last_data = data
                    return f"/stub/{filename}"

            stub = _StubStorage()
            svc = ExportService(storage=stub)  # type: ignore[arg-type]
            url, size = svc._export_pptx_sync(task_id, persisted, "Editorial")
            assert size > 0
            assert stub.last_data is not None

            # Parse the PPTX and confirm edited text is present.
            prs = Presentation(io.BytesIO(stub.last_data))
            text_blob: list[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                text_blob.append(run.text or "")
            joined = " | ".join(text_blob)
            # Edited title and edited bullets are rendered. The PPTX
            # title-slide renderer uppercases and may split the title
            # across two text boxes, so assert on the uppercase pieces.
            assert "EDITED TITLE" in joined.upper()
            assert "XYZ" in joined.upper()
            assert "EDITED BULLETS" in joined.upper()
            assert "alpha edited" in joined
            # Original text is gone.
            assert "Original Title" not in joined
            assert "Original Bullets" not in joined
        finally:
            await _teardown_app(engine)

    asyncio.run(_go())


# ── 6. PUT on missing / unfinished task ───────────────────────────────────


def test_put_unknown_task_returns_404():
    async def _go():
        Session, engine, _ = await _setup_app()
        try:
            async with _client() as c:
                r = await c.put(
                    "/api/slides/does-not-exist",
                    json={"slides": _edited_deck()},
                )
            assert r.status_code == 404
        finally:
            await _teardown_app(engine)

    asyncio.run(_go())


def test_put_pending_task_returns_409():
    async def _go():
        engine = create_async_engine("sqlite+aiosqlite:///:memory:", future=True)
        async with engine.begin() as conn:
            await conn.run_sync(Base.metadata.create_all)
        Session = async_sessionmaker(engine, class_=AsyncSession, expire_on_commit=False)

        async def _get_db_override() -> AsyncIterator[AsyncSession]:
            async with Session() as s:
                yield s

        app.dependency_overrides[get_db] = _get_db_override
        task_id = "task-pending"
        async with Session() as s:
            s.add(Task(id=task_id, topic="t", slide_count=3, theme="Editorial", status="pending"))
            await s.commit()
        try:
            async with _client() as c:
                r = await c.put(
                    f"/api/slides/{task_id}",
                    json={"slides": _edited_deck()},
                )
            assert r.status_code == 409
        finally:
            await _teardown_app(engine)

    asyncio.run(_go())
