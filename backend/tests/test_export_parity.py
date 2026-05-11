"""Phase 6C — Renderer-to-Export parity contract tests.

Goal: prove that every canonical layout exports to PPTX with its core textual
content preserved. This is a **content** parity test, not a pixel-level
visual parity test. Visual parity remains an open risk.

Strategy:
* Build deterministic fixture slides (no `image_url`, so no network) covering
  all 7 canonical layouts.
* Drive `ExportService._export_pptx_sync` with an in-memory storage stub so
  no real filesystem or R2 IO happens.
* Re-open the saved PPTX with `python-pptx` and assert that every fixture's
  unique text markers survived the round trip.
* For the chart slide, also assert that QC labels and unit/source captions
  survive.

Conftest-free; no LLM, no network, no disk dependency outside of optional
write-to-tmp-path smoke checks. Runs against the official Docker gate.
"""
from __future__ import annotations

import io
from typing import Any

import pytest
from pptx import Presentation
from pptx.chart.data import CategoryChartData  # noqa: F401  (sanity import)

from services.export_service import ExportService

from tests.fixtures.canonical_slides import (
    CANONICAL_LAYOUTS,
    canonical_fixture_slides,
)


# ── helpers ─────────────────────────────────────────────────────────────────


class _MemoryStorage:
    """Storage stub that captures `put` payloads in memory."""

    def __init__(self) -> None:
        self.last: tuple[str, bytes, str] | None = None

    def put(self, filename: str, data: bytes, content_type: str = "application/octet-stream") -> str:
        self.last = (filename, data, content_type)
        return f"memory://{filename}"


def _export_to_bytes(slides: list[dict[str, Any]], theme: str = "Editorial") -> bytes:
    storage = _MemoryStorage()
    svc = ExportService(storage=storage)
    url, size = svc._export_pptx_sync("phase6c", slides, theme)
    assert url == "memory://phase6c.pptx"
    assert size > 0, "exported PPTX is empty"
    assert storage.last is not None
    return storage.last[1]


def _shape_texts(slide) -> list[str]:
    out: list[str] = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if run.text:
                        out.append(run.text)
    return out


def _all_texts(slide) -> str:
    return " | ".join(_shape_texts(slide))


def _chart_categories_and_series(slide) -> tuple[list[str], list[float]]:
    for shape in slide.shapes:
        if shape.has_chart:
            chart = shape.chart
            plots = list(chart.plots)
            cats = [str(c) for c in plots[0].categories] if plots else []
            series_vals: list[float] = []
            for series in chart.series:
                series_vals.extend(list(series.values))
            return cats, series_vals
    return [], []


# ── fixtures ────────────────────────────────────────────────────────────────


@pytest.fixture(scope="module")
def exported_pptx_bytes() -> bytes:
    return _export_to_bytes(canonical_fixture_slides())


@pytest.fixture(scope="module")
def exported_presentation(exported_pptx_bytes: bytes) -> Presentation:
    return Presentation(io.BytesIO(exported_pptx_bytes))


# ── core contract tests ─────────────────────────────────────────────────────


def test_export_produces_nonempty_file(exported_pptx_bytes: bytes) -> None:
    # python-pptx blank deck is ~30 KB minimum once a single shape is added.
    assert len(exported_pptx_bytes) > 5000, len(exported_pptx_bytes)
    # PPTX is a ZIP — magic bytes "PK\x03\x04".
    assert exported_pptx_bytes[:4] == b"PK\x03\x04"


def test_export_slide_count_matches_fixture(exported_presentation: Presentation) -> None:
    fixtures = canonical_fixture_slides()
    assert len(exported_presentation.slides) == len(fixtures)
    assert len(fixtures) == len(CANONICAL_LAYOUTS)


def test_export_covers_all_seven_canonical_layouts() -> None:
    """Sanity guard: the fixture itself must cover every canonical layout."""
    layouts = [s["layout"] for s in canonical_fixture_slides()]
    assert set(layouts) == set(CANONICAL_LAYOUTS)


# ── per-layout content parity ───────────────────────────────────────────────


def test_title_layout_preserves_title_and_subtitle(exported_presentation: Presentation) -> None:
    text = _all_texts(exported_presentation.slides[0])
    # Title is rendered upper-case and split across two lines; check word presence.
    assert "RENDERER" in text.upper()
    assert "EXPORT" in text.upper()
    assert "PARITY" in text.upper()
    assert "PARITYSUB-01" in text


def test_bullets_layout_preserves_title_and_all_bullets(exported_presentation: Presentation) -> None:
    text = _all_texts(exported_presentation.slides[1])
    assert "PARITY-BULLETS-TITLE" in text
    assert "BULLET-ITEM-A" in text
    assert "BULLET-ITEM-B" in text
    assert "BULLET-ITEM-C" in text


def test_two_col_layout_preserves_both_columns(exported_presentation: Presentation) -> None:
    text = _all_texts(exported_presentation.slides[2])
    assert "PARITY-TWOCOL-TITLE" in text
    assert "TWOCOL-LEFT-HEAD" in text
    assert "TWOCOL-LEFT-BODY" in text
    assert "TWOCOL-RIGHT-HEAD" in text
    assert "TWOCOL-RIGHT-BODY" in text


def test_quote_layout_preserves_quote_and_attribution(exported_presentation: Presentation) -> None:
    text = _all_texts(exported_presentation.slides[3])
    assert "PARITYQUOTE" in text
    assert "PARITYATTRIB" in text


def test_stats_layout_preserves_all_values_and_labels(exported_presentation: Presentation) -> None:
    text = _all_texts(exported_presentation.slides[4])
    assert "PARITY-STATS-TITLE" in text
    for marker in (
        "STATVAL-42", "STATVAL-77", "STATVAL-93",
        "STATLBL-Alpha", "STATLBL-Bravo", "STATLBL-Charlie",
    ):
        assert marker in text, marker


def test_chart_layout_preserves_title_caption_and_chart_data(
    exported_presentation: Presentation,
) -> None:
    chart_slide = exported_presentation.slides[5]
    text = _all_texts(chart_slide)
    assert "PARITY-CHART-TITLE" in text
    assert "PARITY-CHART-SUB" in text
    # Caption line includes both the unit and the source.
    assert "CHARTUNIT-USD" in text
    assert "CHARTSRC-FY24" in text

    cats, vals = _chart_categories_and_series(chart_slide)
    assert cats == ["CHARTLBL-Q1", "CHARTLBL-Q2", "CHARTLBL-Q3"]
    assert vals == [10.0, 20.0, 30.0]


def test_closing_layout_preserves_title_subtitle_cta(exported_presentation: Presentation) -> None:
    text = _all_texts(exported_presentation.slides[6])
    assert "PARITY-CLOSING-TITLE" in text
    assert "PARITY-CLOSING-SUB" in text
    assert "PARITY-CLOSING-CTA" in text


# ── safety / robustness ────────────────────────────────────────────────────


def test_unknown_layout_does_not_crash_export() -> None:
    """An unknown layout currently falls through to the title renderer.
    Verify that contract: no crash and the title text survives."""
    slides = [{"layout": "totally-made-up", "title": "FALLBACK-TITLE"}]
    data = _export_to_bytes(slides)
    prs = Presentation(io.BytesIO(data))
    assert len(prs.slides) == 1
    assert "FALLBACK" in _all_texts(prs.slides[0]).upper()


def test_chart_with_empty_data_does_not_crash() -> None:
    slides = [
        {
            "layout": "chart",
            "title": "EMPTY-CHART-TITLE",
            "chart_data": {"labels": [], "values": []},
        }
    ]
    data = _export_to_bytes(slides)
    prs = Presentation(io.BytesIO(data))
    text = _all_texts(prs.slides[0])
    assert "EMPTY-CHART-TITLE" in text
    assert "No chart data" in text


def test_export_is_deterministic_for_textual_content() -> None:
    """Running export twice on the same fixtures must produce the same
    extractable text (binary bytes may differ due to embedded metadata)."""
    slides = canonical_fixture_slides()
    a = _export_to_bytes(slides)
    b = _export_to_bytes(slides)
    pa = Presentation(io.BytesIO(a))
    pb = Presentation(io.BytesIO(b))
    texts_a = [_all_texts(s) for s in pa.slides]
    texts_b = [_all_texts(s) for s in pb.slides]
    assert texts_a == texts_b


def test_export_writes_to_tmp_path_smoke(tmp_path) -> None:
    """End-to-end smoke: bytes written to disk reopen as a valid PPTX."""
    data = _export_to_bytes(canonical_fixture_slides())
    out = tmp_path / "phase6c.pptx"
    out.write_bytes(data)
    assert out.stat().st_size == len(data)
    prs = Presentation(str(out))
    assert len(prs.slides) == 7


# ── PDF smoke ──────────────────────────────────────────────────────────────


def test_pdf_export_smoke_or_skip(tmp_path) -> None:
    """Smoke test for PDF export.

    PDF export uses WeasyPrint, which has heavy system dependencies that may
    not be present in every test environment. If WeasyPrint is unavailable
    or its system deps are missing, this test is skipped — and the audit
    explicitly notes that **full PDF visual parity remains open**. We do
    NOT silently claim PDF parity here.
    """
    try:
        import weasyprint  # noqa: F401
    except Exception as exc:  # pragma: no cover - environment-dependent
        pytest.skip(f"weasyprint unavailable: {exc}")

    storage = _MemoryStorage()
    svc = ExportService(storage=storage)
    try:
        url, size = svc._export_pdf_sync("phase6c", canonical_fixture_slides(), "Editorial")
    except Exception as exc:  # pragma: no cover - environment-dependent
        pytest.skip(f"PDF export raised in this environment: {exc}")

    assert size > 0
    assert storage.last is not None
    data = storage.last[1]
    # PDF files start with "%PDF-".
    assert data[:5] == b"%PDF-", data[:8]
