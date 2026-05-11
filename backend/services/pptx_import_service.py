"""Phase 6S — PPTX ingestion.

Reads a user-uploaded ``.pptx`` archive via :mod:`python-pptx` and converts
each slide into a canonical NEXUS slide payload that passes
:func:`agent.slide_schema.validate_deck`. The output is intentionally
conservative: every imported slide becomes either a ``title`` slide
(first slide only, with a non-empty title) or a ``bullets`` slide. We
deliberately do NOT try to reverse-engineer ``two-col``, ``stats``,
``chart``, or ``quote`` layouts from arbitrary PPTX structure — that path
is brittle and out of scope for Phase 6S. Users who want richer layouts
can change the layout in the deck workspace after import.

The function returns a list of canonical slides plus a ``meta`` dict; the
caller is responsible for persistence.
"""

from __future__ import annotations

import io
import logging
from dataclasses import dataclass
from typing import Any

logger = logging.getLogger("nexus.services.pptx_import")


class PPTXImportError(Exception):
    """Raised when the uploaded file is not a usable .pptx archive."""

    def __init__(self, code: str, message: str) -> None:
        super().__init__(message)
        self.code = code
        self.message = message


@dataclass(frozen=True)
class ImportedDeck:
    """Output of :func:`import_pptx_bytes`."""

    title: str
    slides: list[dict[str, Any]]
    source_filename: str
    source_slide_count: int


# Conservative caps so a single hostile PPTX cannot blow up the schema.
_MAX_TITLE_LEN = 200
_MAX_SUBTITLE_LEN = 240
_MAX_BULLET_LEN = 240
_MAX_BULLETS = 4
_MAX_SLIDES = 60  # hard ceiling regardless of source deck size


def _clamp(text: str, limit: int) -> str:
    text = (text or "").strip().replace("\r", " ").replace("\n", " ")
    text = " ".join(text.split())
    if len(text) > limit:
        text = text[: max(1, limit - 1)].rstrip() + "…"
    return text


def _slide_title(slide: Any) -> str:
    """Extract the slide title text, if any.

    Tries the title placeholder first (most reliable), then falls back to
    the first non-empty text frame on the slide.
    """
    try:
        tph = slide.shapes.title
    except Exception:
        tph = None
    if tph is not None:
        try:
            txt = (tph.text or "").strip()
            if txt:
                return txt
        except Exception:
            pass

    for shape in getattr(slide.shapes, "_shapes", []) or slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        try:
            txt = (shape.text_frame.text or "").strip()
        except Exception:
            continue
        if txt:
            # Take only the first paragraph / first line as the title.
            first_line = txt.splitlines()[0].strip()
            if first_line:
                return first_line
    return ""


def _slide_body_lines(slide: Any, *, exclude: str) -> list[str]:
    """Collect non-title body text as a flat list of trimmed lines."""
    lines: list[str] = []
    title_norm = (exclude or "").strip().lower()

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        # Skip the title placeholder itself.
        try:
            ph = getattr(shape, "placeholder_format", None)
            if ph is not None and ph.type is not None and "title" in str(ph.type).lower():
                continue
        except Exception:
            pass
        try:
            tf_text = shape.text_frame.text or ""
        except Exception:
            continue
        for raw in tf_text.splitlines():
            txt = raw.strip().lstrip("•◦●·–-").strip()
            if not txt:
                continue
            if title_norm and txt.strip().lower() == title_norm:
                continue
            lines.append(txt)
    return lines


def _build_title_slide(title: str, subtitle: str) -> dict[str, Any]:
    return {
        "layout": "title",
        "title": _clamp(title or "Imported deck", _MAX_TITLE_LEN),
        "subtitle": _clamp(subtitle, _MAX_SUBTITLE_LEN),
        "eyebrow": "",
    }


def _build_bullets_slide(title: str, body_lines: list[str]) -> dict[str, Any]:
    bullets: list[str] = []
    for line in body_lines:
        clamped = _clamp(line, _MAX_BULLET_LEN)
        if clamped:
            bullets.append(clamped)
        if len(bullets) >= _MAX_BULLETS:
            break
    if not bullets:
        # Schema requires at least one non-empty bullet.
        bullets = ["(no body text imported from this slide)"]
    return {
        "layout": "bullets",
        "title": _clamp(title or "Untitled slide", _MAX_TITLE_LEN),
        "bullets": bullets,
    }


def import_pptx_bytes(data: bytes, *, filename: str = "import.pptx") -> ImportedDeck:
    """Parse a PPTX byte blob into a canonical NEXUS deck.

    Raises
    ------
    PPTXImportError
        With ``code="empty"`` if the upload has no slides, ``code="corrupt"``
        if python-pptx cannot open the archive, or ``code="too_large"`` /
        ``code="bad_extension"`` for caller-side validation reuse.
    """
    if not data:
        raise PPTXImportError("empty_payload", "Uploaded file is empty.")

    # Defer the import so the dependency only loads when this route is hit.
    try:
        from pptx import Presentation
    except Exception as exc:  # pragma: no cover - install-time issue
        logger.exception("pptx_import.python_pptx_unavailable")
        raise PPTXImportError("dependency", f"python-pptx unavailable: {exc}") from exc

    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as exc:
        logger.info("pptx_import.open_failed", extra={"err": str(exc)[:200]})
        raise PPTXImportError(
            "corrupt", "File is not a valid .pptx archive or is corrupt."
        ) from exc

    raw_slides = list(prs.slides)
    source_count = len(raw_slides)
    if source_count == 0:
        raise PPTXImportError("empty", "PPTX has no slides.")

    if source_count > _MAX_SLIDES:
        logger.info(
            "pptx_import.truncated",
            extra={"source": source_count, "limit": _MAX_SLIDES},
        )
        raw_slides = raw_slides[:_MAX_SLIDES]

    slides: list[dict[str, Any]] = []
    deck_title = ""

    for i, src in enumerate(raw_slides):
        title = _slide_title(src)
        body = _slide_body_lines(src, exclude=title)

        if i == 0:
            # First slide: prefer a title-layout slide if we can produce a
            # non-empty title. Subtitle = first body line, when available.
            if title:
                subtitle = body[0] if body else ""
                slides.append(_build_title_slide(title, subtitle))
                deck_title = title
                continue
            # No usable title on slide 1: fall through to bullets layout
            # so we still emit a schema-valid slide.

        slides.append(_build_bullets_slide(title or f"Slide {i + 1}", body))
        if not deck_title and title:
            deck_title = title

    if not slides:  # pragma: no cover - defensive
        raise PPTXImportError("empty", "No slides could be imported.")

    if not deck_title:
        # Use the source filename (without extension) as a fallback topic.
        base = filename.rsplit("/", 1)[-1].rsplit("\\", 1)[-1]
        if base.lower().endswith(".pptx"):
            base = base[:-5]
        deck_title = _clamp(base or "Imported deck", _MAX_TITLE_LEN)

    return ImportedDeck(
        title=deck_title,
        slides=slides,
        source_filename=filename,
        source_slide_count=source_count,
    )


__all__ = ["ImportedDeck", "PPTXImportError", "import_pptx_bytes"]
