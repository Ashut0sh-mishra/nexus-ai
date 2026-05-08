"""Multi-format context extraction for uploaded files.

Each parser is best-effort: it returns whatever text and structured data it
can pull out, and never raises to the caller. Parsing failures degrade to an
empty-but-valid result so the upload pipeline can still record the file.

Public entrypoint: ``extract(path, file_type) -> ExtractionResult``.
"""

from __future__ import annotations

import json
import logging
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

logger = logging.getLogger("nexus.services.context_extractor")


@dataclass
class ExtractionResult:
    """Normalized output from any parser."""

    text: str = ""
    data: dict[str, Any] = field(default_factory=dict)
    error: str | None = None

    def preview(self, limit: int = 500) -> str:
        return (self.text or "").strip()[:limit]


# --------------------------------------------------------------------------- #
# Public dispatcher
# --------------------------------------------------------------------------- #
def extract(path: str | Path, file_type: str) -> ExtractionResult:
    """Run the parser for ``file_type`` and return a normalized result."""
    p = Path(path)
    ft = (file_type or "").lower().lstrip(".")
    try:
        if ft == "csv":
            return _parse_csv(p)
        if ft in {"xlsx", "xls"}:
            return _parse_xlsx(p)
        if ft == "json":
            return _parse_json(p)
        if ft == "pdf":
            return _parse_pdf(p)
        if ft == "docx":
            return _parse_docx(p)
        if ft == "pptx":
            return _parse_pptx(p)
        if ft in {"txt", "md"}:
            return _parse_text(p)
        return ExtractionResult(error=f"Unsupported file type: {ft}")
    except Exception as exc:  # pragma: no cover — last-resort safety
        logger.warning(
            "context.extract_failed",
            extra={"path": str(p), "file_type": ft, "err": str(exc)},
        )
        # Try a raw text read so the file isn't completely opaque.
        return _parse_text(p, fallback_error=str(exc))


# --------------------------------------------------------------------------- #
# Parsers
# --------------------------------------------------------------------------- #
def _parse_csv(path: Path) -> ExtractionResult:
    try:
        import pandas as pd  # type: ignore
    except ImportError:
        return _parse_text(path, fallback_error="pandas not installed")

    try:
        df = pd.read_csv(path, nrows=5000)
    except Exception as exc:
        return _parse_text(path, fallback_error=f"csv parse failed: {exc}")

    return _from_dataframe(df, path.name)


def _parse_xlsx(path: Path) -> ExtractionResult:
    try:
        import pandas as pd  # type: ignore
    except ImportError:
        return ExtractionResult(error="pandas not installed")

    try:
        # Read every sheet; cap rows per sheet for memory safety.
        sheets = pd.read_excel(path, sheet_name=None, nrows=5000)
    except Exception as exc:
        return ExtractionResult(error=f"xlsx parse failed: {exc}")

    sheet_results: list[dict[str, Any]] = []
    text_parts: list[str] = []
    for name, df in sheets.items():
        sub = _from_dataframe(df, f"{path.name} :: {name}")
        sheet_results.append({"sheet": name, **sub.data})
        text_parts.append(f"# Sheet: {name}\n{sub.text}")

    return ExtractionResult(
        text="\n\n".join(text_parts),
        data={"format": "xlsx", "sheets": sheet_results},
    )


def _from_dataframe(df: Any, label: str) -> ExtractionResult:
    """Common normalization for CSV / single-sheet XLSX."""
    headers = [str(c) for c in df.columns.tolist()]
    rows_total = int(len(df))
    sample_rows = df.head(5).fillna("").astype(str).values.tolist()

    numeric_cols: list[dict[str, Any]] = []
    for col in df.columns:
        try:
            series = df[col].dropna()
            # pandas keeps numeric dtypes; coerce to be safe
            import pandas as pd  # type: ignore

            numeric = pd.to_numeric(series, errors="coerce").dropna()
            if len(numeric) >= 2 and len(numeric) >= 0.5 * len(series):
                numeric_cols.append(
                    {
                        "name": str(col),
                        "min": float(numeric.min()),
                        "max": float(numeric.max()),
                        "avg": float(numeric.mean()),
                        "sum": float(numeric.sum()),
                        "count": int(len(numeric)),
                    }
                )
        except Exception:
            continue

    # Plain-text rendering for LLM consumption.
    lines = ["\t".join(headers)]
    for row in sample_rows:
        lines.append("\t".join(str(c) for c in row))
    text = f"{label}\n" + "\n".join(lines)

    return ExtractionResult(
        text=text,
        data={
            "format": "table",
            "headers": headers,
            "row_count": rows_total,
            "sample_rows": sample_rows,
            "numeric_columns": numeric_cols,
        },
    )


def _parse_json(path: Path) -> ExtractionResult:
    try:
        with path.open("r", encoding="utf-8") as fh:
            obj = json.load(fh)
    except Exception as exc:
        return _parse_text(path, fallback_error=f"json parse failed: {exc}")

    arrays: list[dict[str, Any]] = []
    if isinstance(obj, list):
        arrays.append({"path": "$", "length": len(obj)})
    elif isinstance(obj, dict):
        for k, v in obj.items():
            if isinstance(v, list):
                arrays.append({"path": f"$.{k}", "length": len(v)})

    return ExtractionResult(
        text=json.dumps(obj, indent=2, ensure_ascii=False)[:20000],
        data={"format": "json", "arrays": arrays, "top_level_type": type(obj).__name__},
    )


def _parse_pdf(path: Path) -> ExtractionResult:
    try:
        import pdfplumber  # type: ignore
    except ImportError:
        return ExtractionResult(error="pdfplumber not installed")

    text_parts: list[str] = []
    tables: list[dict[str, Any]] = []
    page_count = 0
    try:
        with pdfplumber.open(path) as pdf:
            page_count = len(pdf.pages)
            for i, page in enumerate(pdf.pages[:50]):  # cap pages
                page_text = page.extract_text() or ""
                if page_text:
                    text_parts.append(page_text)
                try:
                    page_tables = page.extract_tables() or []
                    for t in page_tables:
                        if not t or not t[0]:
                            continue
                        tables.append(
                            {
                                "page": i + 1,
                                "headers": [str(c or "") for c in t[0]],
                                "rows": [
                                    [str(c or "") for c in row] for row in t[1:21]
                                ],
                            }
                        )
                except Exception:
                    pass
    except Exception as exc:
        return ExtractionResult(error=f"pdf parse failed: {exc}")

    return ExtractionResult(
        text="\n\n".join(text_parts),
        data={"format": "pdf", "page_count": page_count, "tables": tables},
    )


def _parse_docx(path: Path) -> ExtractionResult:
    try:
        from docx import Document  # type: ignore
    except ImportError:
        return ExtractionResult(error="python-docx not installed")

    try:
        doc = Document(str(path))
    except Exception as exc:
        return ExtractionResult(error=f"docx parse failed: {exc}")

    headings: list[dict[str, Any]] = []
    text_parts: list[str] = []
    for para in doc.paragraphs:
        s = (para.style.name or "") if para.style else ""
        text = para.text or ""
        if not text.strip():
            continue
        if s.startswith("Heading"):
            try:
                level = int("".join(ch for ch in s if ch.isdigit()) or "1")
            except ValueError:
                level = 1
            headings.append({"level": level, "text": text})
            text_parts.append(("#" * min(level, 6)) + " " + text)
        else:
            text_parts.append(text)

    return ExtractionResult(
        text="\n\n".join(text_parts),
        data={"format": "docx", "headings": headings},
    )


def _parse_pptx(path: Path) -> ExtractionResult:
    try:
        from pptx import Presentation  # type: ignore
    except ImportError:
        return ExtractionResult(error="python-pptx not installed")

    try:
        prs = Presentation(str(path))
    except Exception as exc:
        return ExtractionResult(error=f"pptx parse failed: {exc}")

    slides: list[dict[str, Any]] = []
    text_parts: list[str] = []
    for i, slide in enumerate(prs.slides):
        slide_text_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text or "" for run in para.runs).strip()
                    if line:
                        slide_text_parts.append(line)
        slide_text = "\n".join(slide_text_parts)
        slides.append({"slide_number": i + 1, "text": slide_text})
        if slide_text:
            text_parts.append(f"## Slide {i + 1}\n{slide_text}")

    return ExtractionResult(
        text="\n\n".join(text_parts),
        data={"format": "pptx", "slide_count": len(slides), "slides": slides},
    )


def _parse_text(path: Path, fallback_error: str | None = None) -> ExtractionResult:
    try:
        text = path.read_text(encoding="utf-8", errors="replace")
    except Exception as exc:
        return ExtractionResult(error=fallback_error or f"text read failed: {exc}")
    return ExtractionResult(
        text=text,
        data={"format": "text", "length": len(text)},
        error=fallback_error,
    )
