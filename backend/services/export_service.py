"""PPTX + PDF export, with R2 or local-filesystem storage."""

from __future__ import annotations

import asyncio
import io
import logging
from pathlib import Path
from typing import Any, Iterable

import httpx
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

from config import settings
from services.storage_service import StorageService
from agent.themes_registry import get_theme

logger = logging.getLogger("nexus.services.export")

# 16:9 master dimensions in EMU
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _palette_for(theme: str | None) -> dict[str, str]:
    """Phase 6O — resolve theme tokens via the registry.

    Returns hex strings without the leading ``#`` so the existing PPTX
    helpers (``_hex_to_rgb``, background fills, etc.) keep working
    unchanged. Legacy display names (``Editorial``, ``Pixel``, ``Vellum``,
    ``Dossier``, ``light-pro``) all resolve to the registry now.
    """

    t = get_theme(theme)
    c = t.colors
    strip = lambda v: v.lstrip("#")  # noqa: E731
    return {
        "bg": strip(c.get("bg", "#0F0F14")),
        "text": strip(c.get("text", "#F5F5F7")),
        "muted": strip(c.get("muted", "#9A9AA5")),
        "accent": strip(c.get("accent", "#A78BFA")),
        # ``chart_palette`` is included as a comma-joined hex string so the
        # QuickChart helper can split and reuse it without reaching back
        # into the registry from inside a static method.
        "chart_palette": ",".join(strip(x) for x in t.chart_palette),
    }


# Backwards-compatible legacy palette dict, now derived from the registry
# instead of being hard-coded. Keys preserve the legacy display names that
# already exist in ``Task.theme`` rows and in the React renderer.
THEMES: dict[str, dict[str, str]] = {
    name: _palette_for(name)
    for name in (
        "light-pro", "Editorial", "Pixel", "Vellum", "Dossier",
        "Whiteboard", "Sketch", "Glamour", "Amber", "Arctic",
        "Cerulean", "Cobalt", "Emerald", "Basalt", "Mist",
        "Onyx", "Sand", "Neon", "Linen",
    )
}


def _hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class ExportService:
    def __init__(self, storage: StorageService | None = None) -> None:
        self.storage = storage or StorageService()

    # ── public ────────────────────────────────────────────────────────────
    async def export_pptx(
        self, task_id: str, slides: list[dict[str, Any]], theme: str
    ) -> tuple[str, int]:
        return await asyncio.to_thread(self._export_pptx_sync, task_id, slides, theme)

    async def export_pdf(
        self, task_id: str, slides: list[dict[str, Any]], theme: str
    ) -> tuple[str, int]:
        return await asyncio.to_thread(self._export_pdf_sync, task_id, slides, theme)

    # ── pptx ──────────────────────────────────────────────────────────────
    def _export_pptx_sync(
        self, task_id: str, slides: list[dict[str, Any]], theme: str
    ) -> tuple[str, int]:
        palette = THEMES.get(theme, THEMES["Editorial"])
        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H

        # Phase 6AL-Export: prefetch every distinct image URL in parallel.
        # Pre-6AL this was sequential inside ``_render_slide``: each image
        # could spend up to 25s x 3 retries = ~80s, so an 8-slide deck with
        # images on 4-6 slides could take 5+ minutes before the PPTX even
        # started rendering. The parallel prefetch caps total image-wait at
        # ~12s regardless of how many slides have images. A missing image
        # is not fatal: ``_render_slide`` already handles ``None`` bytes.
        image_cache = self._prefetch_images(slides)

        for slide in slides:
            self._render_slide(prs, slide, palette, image_cache=image_cache)

        buf = io.BytesIO()
        prs.save(buf)
        data = buf.getvalue()
        filename = f"{task_id}.pptx"
        url = self.storage.put(
            filename,
            data,
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        logger.info("export.pptx_ok", extra={"task_id": task_id, "size": len(data)})
        return url, len(data)

    def _render_slide(
        self,
        prs: Presentation,
        slide: dict[str, Any],
        palette: dict[str, str],
        *,
        image_cache: dict[str, bytes | None] | None = None,
    ) -> None:
        layout = (slide.get("layout") or "title").lower()
        s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        self._add_background(s, palette["bg"])

        # Hero image (best-effort; never blocks export). When ``image_cache``
        # is provided (typical PPTX path) we look up the prefetched bytes
        # in O(1). Fallback to a direct fetch keeps the method standalone
        # for callers that bypass ``_export_pptx_sync``.
        url = slide.get("image_url")
        if image_cache is not None and url is not None:
            image_bytes = image_cache.get(url)
        else:
            image_bytes = self._fetch_image(url)
        if image_bytes:
            if layout == "closing":
                # Full-bleed with semi-transparent dark scrim.
                s.shapes.add_picture(
                    io.BytesIO(image_bytes), 0, 0, width=SLIDE_W, height=SLIDE_H
                )
                scrim = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
                scrim.line.fill.background()
                scrim.fill.solid()
                scrim.fill.fore_color.rgb = _hex_to_rgb(palette["bg"])
                scrim.fill.transparency = 0.35
                scrim.shadow.inherit = False
            elif layout in ("bullets", "two-col"):
                # Right-side panel covering ~38% of the width.
                pic_w = Inches(5.0)
                s.shapes.add_picture(
                    io.BytesIO(image_bytes),
                    SLIDE_W - pic_w,
                    0,
                    width=pic_w,
                    height=SLIDE_H,
                )
            # NOTE: title layout owns its own image placement (right panel)
            # via _render_title's split-screen composition.

        if layout == "title":
            self._render_title(s, slide, palette, image_bytes=image_bytes)
        elif layout == "bullets":
            self._render_bullets(s, slide, palette, has_image=bool(image_bytes))
        elif layout == "two-col":
            self._render_two_col(s, slide, palette, has_image=bool(image_bytes))
        elif layout == "quote":
            self._render_quote(s, slide, palette)
        elif layout == "stats":
            self._render_stats(s, slide, palette)
        elif layout == "chart":
            self._render_chart(s, slide, palette)
        elif layout == "closing":
            self._render_closing(s, slide, palette)
        elif layout == "bigstat":
            # Phase 6AA — degrade to the stats renderer by promoting the
            # single value/label pair into a one-item ``stats`` payload.
            # Never mutates the source slide.
            stats_proxy = {
                **slide,
                "stats": [
                    {
                        "value": str(slide.get("value") or "—"),
                        "label": str(slide.get("label") or ""),
                    }
                ],
            }
            self._render_stats(s, stats_proxy, palette)
        elif layout == "section_divider":
            # Phase 6AA — degrade to the title renderer; eyebrow/subtitle
            # already align with the title-slide field shape.
            self._render_title(s, slide, palette)
        elif layout == "timeline":
            # Phase 6AC — degrade to bullets by formatting each event as
            # "DATE — LABEL". The slide's title is preserved; bullet
            # ordering matches event ordering.
            events = slide.get("events") or []
            bullets_proxy = {
                **slide,
                "bullets": [
                    f"{(e.get('date') or '').strip()} — {(e.get('label') or '').strip()}".strip(" —")
                    for e in events
                    if isinstance(e, dict)
                ][:4],
            }
            self._render_bullets(s, bullets_proxy, palette, has_image=bool(image_bytes))
        elif layout == "comparison":
            # Phase 6AC — degrade to two-col by promoting left/right
            # blocks into the columns array.
            left = slide.get("left") or {}
            right = slide.get("right") or {}
            two_col_proxy = {
                **slide,
                "columns": [
                    {
                        "heading": str(left.get("heading") or "").strip(),
                        "body": str(left.get("body") or "").strip(),
                    },
                    {
                        "heading": str(right.get("heading") or "").strip(),
                        "body": str(right.get("body") or "").strip(),
                    },
                ],
            }
            self._render_two_col(s, two_col_proxy, palette, has_image=bool(image_bytes))
        else:
            self._render_title(s, slide, palette)

    # Per-image budget for PPTX exports.
    # Phase 6AL-Export: tightened from 25s x 3 retries (worst-case ~80s)
    # to 10s x 2 attempts (worst-case ~12s). A missing image is fine; a
    # hung export is not. The Pollinations 429 case still gets one retry.
    _IMAGE_TIMEOUT_S = 10.0
    _IMAGE_RETRY_DELAYS_S = (0.0, 1.5)

    @classmethod
    def _fetch_image(cls, url: str | None) -> bytes | None:
        if not url:
            return None
        for attempt, delay in enumerate(cls._IMAGE_RETRY_DELAYS_S):
            if delay:
                import time as _t

                _t.sleep(delay)
            try:
                with httpx.Client(
                    timeout=cls._IMAGE_TIMEOUT_S, follow_redirects=True
                ) as client:
                    r = client.get(url)
                if r.status_code == 429 and attempt < len(cls._IMAGE_RETRY_DELAYS_S) - 1:
                    continue
                r.raise_for_status()
                ct = r.headers.get("content-type", "")
                if not ct.startswith("image/"):
                    return None
                return r.content
            except Exception as exc:
                if attempt == len(cls._IMAGE_RETRY_DELAYS_S) - 1:
                    logger.warning(
                        "export.image_fetch_failed",
                        extra={"err": str(exc), "url": url[:120]},
                    )
        return None

    @classmethod
    def _prefetch_images(
        cls, slides: list[dict[str, Any]]
    ) -> dict[str, bytes | None]:
        """Fetch every distinct ``image_url`` in parallel.

        Returns a ``{url: bytes-or-None}`` map. Total wall-clock time is
        bounded by ``_IMAGE_TIMEOUT_S`` x retry count regardless of how
        many slides have images.
        """
        from concurrent.futures import ThreadPoolExecutor

        urls: list[str] = []
        seen: set[str] = set()
        for slide in slides:
            u = slide.get("image_url") if isinstance(slide, dict) else None
            if isinstance(u, str) and u and u not in seen:
                seen.add(u)
                urls.append(u)
        if not urls:
            return {}

        cache: dict[str, bytes | None] = {}
        # Cap concurrency so we do not hammer Pollinations into a 429 storm.
        max_workers = min(8, len(urls))
        with ThreadPoolExecutor(max_workers=max_workers) as ex:
            for u, data in zip(urls, ex.map(cls._fetch_image, urls)):
                cache[u] = data
        ok = sum(1 for v in cache.values() if v)
        logger.info(
            "export.images_prefetched",
            extra={"requested": len(urls), "ok": ok, "missing": len(urls) - ok},
        )
        return cache

    @staticmethod
    def _quickchart_url(slide: dict[str, Any], p: dict[str, str]) -> str:
        """Build a QuickChart.io URL that renders a Chart.js config to PNG."""
        import json as _json
        from urllib.parse import quote as _quote

        cd = slide.get("chart_data") or {}
        labels = cd.get("labels") or []
        values = cd.get("values") or []
        unit = cd.get("unit") or ""
        chart_type = (slide.get("chart_type") or "bar").lower()
        if chart_type not in {"bar", "line", "doughnut"}:
            chart_type = "bar"

        accent = f"#{p['accent']}"
        muted = f"#{p['muted']}"
        is_light = p.get("bg", "").upper().startswith("FFF") or p.get("bg", "").upper().startswith("FAF")
        grid_color = "rgba(0,0,0,0.08)" if is_light else "rgba(255,255,255,0.10)"
        tick_color = muted

        if chart_type == "doughnut":
            chart_palette_raw = p.get("chart_palette") or ""
            chart_palette = [
                f"#{h}" for h in chart_palette_raw.split(",") if h
            ] or [accent]
            bg_colors = chart_palette[: len(labels)] or [accent]
            datasets = [
                {
                    "data": values,
                    "backgroundColor": bg_colors,
                    "borderWidth": 0,
                }
            ]
        else:
            datasets = [
                {
                    "label": f"Value ({unit})" if unit else "Value",
                    "data": values,
                    "backgroundColor": accent + "CC",
                    "borderColor": accent,
                    "borderWidth": 2,
                    "fill": chart_type == "line",
                    "tension": 0.35,
                }
            ]

        scales = {} if chart_type == "doughnut" else {
            "x": {"grid": {"color": grid_color}, "ticks": {"color": tick_color}},
            "y": {
                "grid": {"color": grid_color},
                "ticks": {"color": tick_color},
                "beginAtZero": True,
            },
        }

        config = {
            "type": chart_type,
            "data": {"labels": labels, "datasets": datasets},
            "options": {
                "plugins": {"legend": {"display": chart_type == "doughnut"}},
                "scales": scales,
            },
        }
        encoded = _quote(_json.dumps(config, ensure_ascii=False), safe="")
        # White background so the PNG looks correct on light themes; transparent
        # would be ideal but QuickChart uses white by default.
        return (
            f"https://quickchart.io/chart?w=900&h=420&bkg=transparent&c={encoded}"
        )

    @staticmethod
    def _add_background(slide, hex_color: str) -> None:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
        shape.line.fill.background()
        shape.fill.solid()
        shape.fill.fore_color.rgb = _hex_to_rgb(hex_color)
        shape.shadow.inherit = False

    @staticmethod
    def _add_text(
        slide,
        text: str,
        left: int,
        top: int,
        width: int,
        height: int,
        *,
        size: int,
        color: str,
        bold: bool = False,
        align: str = "left",
        italic: bool = False,
    ):
        from pptx.enum.text import PP_ALIGN

        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_right = Pt(0)
        tf.margin_top = tf.margin_bottom = Pt(0)
        tf.text = text or ""
        for para in tf.paragraphs:
            para.alignment = {
                "left": PP_ALIGN.LEFT,
                "center": PP_ALIGN.CENTER,
                "right": PP_ALIGN.RIGHT,
            }.get(align, PP_ALIGN.LEFT)
            for run in para.runs:
                run.font.size = Pt(size)
                run.font.bold = bold
                run.font.italic = italic
                run.font.color.rgb = _hex_to_rgb(color)
                run.font.name = "Inter"
        return tb

    def _render_title(
        self,
        slide,
        data: dict[str, Any],
        p: dict[str, str],
        *,
        image_bytes: bytes | None = None,
    ) -> None:
        # Phase 6AL-Visuals: cinematic full-bleed cover.
        # Pre-6AL composition was a 58/42 split with a solid accent panel
        # and a sun-glyph disc on the right plus a "POWERED BY NEXUS"
        # footer band \u2014 the single biggest source of "PowerPoint template"
        # energy in the export. We now render the image full-bleed (when
        # available), drape a vertical gradient scrim, and place an
        # editorial title block in the bottom-left. No disc, no accent
        # panel, no brand bar.
        if image_bytes:
            # Full-bleed image
            slide.shapes.add_picture(
                io.BytesIO(image_bytes),
                Inches(0),
                Inches(0),
                width=SLIDE_W,
                height=SLIDE_H,
            )
            # Dark scrim across the bottom 65% so the title is readable
            # regardless of the underlying image.
            scrim_h = Inches(5.0)
            scrim = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(0),
                SLIDE_H - scrim_h,
                SLIDE_W,
                scrim_h,
            )
            scrim.line.fill.background()
            scrim.fill.solid()
            scrim.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
            scrim.fill.transparency = 0.30
            scrim.shadow.inherit = False
            title_color = "#FFFFFF"
            subtitle_color = "#E5E7EB"
            eyebrow_color = "#FFFFFF"
            rule_color = "#FFFFFF"
        else:
            # No image: solid background, accent rule under the title.
            title_color = p["text"]
            subtitle_color = p["muted"]
            eyebrow_color = p["accent"]
            rule_color = p["accent"]

        # Eyebrow (top-left, widely tracked, all caps).
        eyebrow = (data.get("eyebrow") or "Presentation").upper()
        self._add_text(
            slide, eyebrow,
            Inches(0.9), Inches(0.8), Inches(11.5), Inches(0.4),
            size=11, color=eyebrow_color, bold=True,
        )

        # Title block bottom-left, single line for editorial drama. We do
        # NOT split the title into two-colored lines anymore \u2014 that was
        # template energy. python-pptx will wrap if the title overflows.
        full_title = (data.get("title") or "").strip()
        self._add_text(
            slide, full_title,
            Inches(0.9), Inches(4.6), Inches(11.5), Inches(2.0),
            size=60, color=title_color, bold=True,
        )

        if data.get("subtitle"):
            self._add_text(
                slide, data["subtitle"],
                Inches(0.9), Inches(6.6), Inches(11.5), Inches(0.6),
                size=16, color=subtitle_color,
            )

        # Thin accent rule under the title block (editorial signature).
        rule = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(6.45), Inches(1.3), Inches(0.02)
        )
        rule.line.fill.background()
        rule.fill.solid()
        rule.fill.fore_color.rgb = _hex_to_rgb(rule_color)
        rule.shadow.inherit = False

    def _render_bullets(
        self,
        slide,
        data: dict[str, Any],
        p: dict[str, str],
        *,
        has_image: bool = False,
    ) -> None:
        text_w = Inches(7.0) if has_image else Inches(11.7)
        bullet_text_w = Inches(6.5) if has_image else Inches(11.3)
        self._add_text(
            slide, data.get("title", ""),
            Inches(0.8), Inches(0.6), text_w, Inches(1.0),
            size=32, color=p["text"], bold=True,
        )
        bullets: Iterable[str] = data.get("bullets") or []
        y = 2.0
        for b in bullets:
            self._add_text(
                slide, "•",
                Inches(0.8), Inches(y), Inches(0.4), Inches(0.6),
                size=20, color=p["accent"], bold=True,
            )
            self._add_text(
                slide, str(b),
                Inches(1.3), Inches(y), bullet_text_w, Inches(0.9),
                size=18, color=p["text"],
            )
            y += 1.05

    def _render_two_col(
        self,
        slide,
        data: dict[str, Any],
        p: dict[str, str],
        *,
        has_image: bool = False,
    ) -> None:
        title_w = Inches(7.0) if has_image else Inches(11.7)
        self._add_text(
            slide, data.get("title", ""),
            Inches(0.8), Inches(0.6), title_w, Inches(1.0),
            size=32, color=p["text"], bold=True,
        )
        cols = (data.get("columns") or [])[:2]
        if has_image:
            col_w = Inches(7.0)
            for i, c in enumerate(cols):
                y = 2.0 + i * 2.6
                self._add_text(
                    slide, (c.get("heading") or "").strip(),
                    Inches(0.8), Inches(y), col_w, Inches(0.6),
                    size=18, color=p["accent"], bold=True,
                )
                self._add_text(
                    slide, (c.get("body") or "").strip(),
                    Inches(0.8), Inches(y + 0.7), col_w, Inches(1.8),
                    size=15, color=p["text"],
                )
            return
        col_w = Inches(5.6)
        for i, c in enumerate(cols):
            x = Inches(0.8 + i * 6.0)
            self._add_text(
                slide, (c.get("heading") or "").strip(),
                x, Inches(2.0), col_w, Inches(0.6),
                size=18, color=p["accent"], bold=True,
            )
            self._add_text(
                slide, (c.get("body") or "").strip(),
                x, Inches(2.7), col_w, Inches(4.2),
                size=15, color=p["text"],
            )

    def _render_quote(self, slide, data: dict[str, Any], p: dict[str, str]) -> None:
        self._add_text(
            slide, "“",
            Inches(1), Inches(1.6), Inches(11.3), Inches(1.5),
            size=72, color=p["accent"], bold=True, align="center",
        )
        self._add_text(
            slide, data.get("quote") or data.get("title", ""),
            Inches(1.5), Inches(3.0), Inches(10.3), Inches(2.4),
            size=28, color=p["text"], italic=True, align="center",
        )
        if data.get("attribution"):
            self._add_text(
                slide, f"— {data['attribution']}",
                Inches(1.5), Inches(5.6), Inches(10.3), Inches(0.6),
                size=14, color=p["muted"], align="center",
            )

    def _render_stats(self, slide, data: dict[str, Any], p: dict[str, str]) -> None:
        self._add_text(
            slide, data.get("title", ""),
            Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0),
            size=32, color=p["text"], bold=True,
        )
        stats = (data.get("stats") or [])[:3]
        gap = 0.4
        col_w = (13.333 - 1.6 - 2 * gap) / 3
        for i, s in enumerate(stats):
            x = Inches(0.8 + i * (col_w + gap))
            self._add_text(
                slide, str(s.get("value", "")),
                x, Inches(2.6), Inches(col_w), Inches(1.6),
                size=56, color=p["accent"], bold=True, align="center",
            )
            self._add_text(
                slide, str(s.get("label", "")),
                x, Inches(4.4), Inches(col_w), Inches(0.8),
                size=14, color=p["muted"], align="center",
            )

    def _render_chart(self, slide, data: dict[str, Any], p: dict[str, str]) -> None:
        # Title + optional subtitle
        self._add_text(
            slide, data.get("title", ""),
            Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8),
            size=30, color=p["text"], bold=True,
        )
        if data.get("subtitle"):
            self._add_text(
                slide, data["subtitle"],
                Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5),
                size=14, color=p["muted"],
            )

        cd = data.get("chart_data") or {}
        labels = cd.get("labels") or []
        values = cd.get("values") or []
        unit = cd.get("unit") or ""
        source = cd.get("source") or ""

        if not labels or not values:
            self._add_text(
                slide, "No chart data",
                Inches(0.8), Inches(3.5), Inches(11.7), Inches(0.6),
                size=14, color=p["muted"], align="center",
            )
            return

        chart_type_raw = (data.get("chart_type") or "bar").lower()
        xl_type = {
            "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
        }.get(chart_type_raw, XL_CHART_TYPE.COLUMN_CLUSTERED)

        chart_data = CategoryChartData()
        chart_data.categories = [str(x) for x in labels]
        series_label = f"Value ({unit})" if unit else "Value"
        chart_data.add_series(series_label, [float(v) for v in values])

        # Chart frame: leaves room for source line at the bottom.
        chart_x = Inches(0.8)
        chart_y = Inches(2.0)
        chart_w = Inches(11.7)
        chart_h = Inches(4.6)
        graphic = slide.shapes.add_chart(
            xl_type, chart_x, chart_y, chart_w, chart_h, chart_data
        )
        chart = graphic.chart

        # Style: hide legend for single-series bar/line, keep for doughnut.
        if chart_type_raw == "doughnut":
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False
        else:
            chart.has_legend = False

        # Color the series with the theme accent.
        try:
            series = chart.series[0]
            fill = series.format.fill
            fill.solid()
            fill.fore_color.rgb = _hex_to_rgb(p["accent"])
            line = series.format.line
            line.color.rgb = _hex_to_rgb(p["accent"])
        except Exception:  # pragma: no cover - defensive
            pass

        # Source / unit caption.
        caption_parts = []
        if unit:
            caption_parts.append(f"Units: {unit}")
        if source:
            caption_parts.append(f"Source: {source}")
        if caption_parts:
            self._add_text(
                slide, "    \u2022    ".join(caption_parts),
                Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.4),
                size=11, color=p["muted"],
            )

    def _render_closing(self, slide, data: dict[str, Any], p: dict[str, str]) -> None:
        self._add_text(
            slide, data.get("title", "Thank you"),
            Inches(1), Inches(2.6), Inches(11.3), Inches(1.6),
            size=44, color=p["text"], bold=True, align="center",
        )
        if data.get("subtitle"):
            self._add_text(
                slide, data["subtitle"],
                Inches(1.5), Inches(4.2), Inches(10.3), Inches(1.0),
                size=18, color=p["muted"], align="center",
            )
        if data.get("cta"):
            self._add_text(
                slide, data["cta"],
                Inches(4.5), Inches(5.4), Inches(4.3), Inches(0.7),
                size=18, color=p["bg"], bold=True, align="center",
            )

    # ── pdf ───────────────────────────────────────────────────────────────
    def _export_pdf_sync(
        self, task_id: str, slides: list[dict[str, Any]], theme: str
    ) -> tuple[str, int]:
        try:
            from weasyprint import HTML
        except Exception as exc:  # pragma: no cover
            raise RuntimeError(
                "weasyprint is not installed or its system deps are missing"
            ) from exc

        palette = THEMES.get(theme, THEMES["Editorial"])
        html = self._slides_to_html(slides, palette)
        pdf_bytes = HTML(string=html).write_pdf()
        filename = f"{task_id}.pdf"
        url = self.storage.put(filename, pdf_bytes, content_type="application/pdf")
        logger.info("export.pdf_ok", extra={"task_id": task_id, "size": len(pdf_bytes)})
        return url, len(pdf_bytes)

    def _slides_to_html(self, slides: list[dict[str, Any]], p: dict[str, str]) -> str:
        from html import escape

        bg, text, muted, accent = p["bg"], p["text"], p["muted"], p["accent"]
        css = f"""
        @page {{ size: 13.333in 7.5in; margin: 0; }}
        body {{ margin:0; font-family:Inter,system-ui,sans-serif; }}
        .slide {{ width:13.333in; height:7.5in; box-sizing:border-box;
                  padding:0.8in; page-break-after:always;
                  background:#{bg}; color:#{text}; position:relative; }}
        .slide:last-child {{ page-break-after:auto; }}
        .eyebrow {{ font-size:11pt; letter-spacing:0.2em; color:#{accent};
                   text-transform:uppercase; text-align:center; margin-bottom:0.3in; }}
        h1.title {{ font-size:48pt; font-weight:600; text-align:center;
                   line-height:1.05; margin:0; }}
        .subtitle {{ font-size:18pt; color:#{muted}; text-align:center;
                    margin-top:0.4in; }}
        h2 {{ font-size:30pt; font-weight:600; margin:0 0 0.4in 0; }}
        ul {{ list-style:none; padding:0; margin:0; }}
        li {{ font-size:18pt; margin:0.3in 0; padding-left:0.4in; position:relative; }}
        li::before {{ content:"•"; color:#{accent}; position:absolute;
                     left:0; font-weight:700; }}
        .cols {{ display:flex; gap:0.4in; }}
        .col {{ flex:1; border:1px solid #{accent}55; border-radius:14px; padding:0.4in; }}
        .col h3 {{ color:#{accent}; font-size:18pt; margin:0 0 0.2in 0; }}
        .col p {{ font-size:14pt; color:#{text}; margin:0; line-height:1.5; }}
        .quote {{ display:flex; flex-direction:column; align-items:center;
                 justify-content:center; height:100%; text-align:center; }}
        .quote .mark {{ font-size:80pt; color:#{accent}; line-height:0.6;
                       margin-bottom:0.4in; }}
        .quote blockquote {{ font-size:26pt; font-style:italic; max-width:9in;
                            margin:0; }}
        .quote .attr {{ font-size:13pt; color:#{muted}; margin-top:0.4in;
                       letter-spacing:0.15em; text-transform:uppercase; }}
        .stats {{ display:flex; gap:0.3in; margin-top:0.4in; }}
        .stat {{ flex:1; border:1px solid #{accent}55; border-radius:14px;
                padding:0.5in; text-align:center; }}
        .stat .v {{ font-size:50pt; font-weight:600; color:#{accent}; }}
        .stat .l {{ font-size:13pt; color:#{muted}; margin-top:0.15in;
                   letter-spacing:0.1em; text-transform:uppercase; }}
        .closing {{ display:flex; flex-direction:column; align-items:center;
                   justify-content:center; height:100%; text-align:center; }}
        .cta {{ display:inline-block; margin-top:0.5in; padding:0.2in 0.6in;
               background:#{accent}; color:#{bg}; border-radius:14px;
               font-weight:600; font-size:16pt; }}
        .chart-wrap {{ width:100%; height:4.6in; display:flex; align-items:center;
                      justify-content:center; margin-top:0.2in; }}
        .chart-wrap img {{ max-width:100%; max-height:100%; }}
        .chart-caption {{ margin-top:0.2in; font-size:11pt; color:#{muted}; }}
        """
        out = [f"<html><head><style>{css}</style></head><body>"]
        for s in slides:
            layout = (s.get("layout") or "title").lower()
            out.append('<section class="slide">')
            if layout == "title":
                out.append(f'<div class="eyebrow">{escape(s.get("eyebrow") or "Presentation")}</div>')
                out.append(f'<h1 class="title">{escape(s.get("title", ""))}</h1>')
                if s.get("subtitle"):
                    out.append(f'<div class="subtitle">{escape(s["subtitle"])}</div>')
            elif layout == "bullets":
                out.append(f'<h2>{escape(s.get("title", ""))}</h2><ul>')
                for b in (s.get("bullets") or [])[:6]:
                    out.append(f"<li>{escape(str(b))}</li>")
                out.append("</ul>")
            elif layout == "two-col":
                out.append(f'<h2>{escape(s.get("title", ""))}</h2><div class="cols">')
                for c in (s.get("columns") or [])[:2]:
                    out.append(
                        f'<div class="col"><h3>{escape(c.get("heading", ""))}</h3>'
                        f'<p>{escape(c.get("body", ""))}</p></div>'
                    )
                out.append("</div>")
            elif layout == "quote":
                out.append('<div class="quote"><div class="mark">“</div>')
                out.append(f'<blockquote>{escape(s.get("quote", s.get("title", "")))}</blockquote>')
                if s.get("attribution"):
                    out.append(f'<div class="attr">— {escape(s["attribution"])}</div>')
                out.append("</div>")
            elif layout == "stats":
                out.append(f'<h2>{escape(s.get("title", ""))}</h2><div class="stats">')
                for st in (s.get("stats") or [])[:3]:
                    out.append(
                        f'<div class="stat"><div class="v">{escape(str(st.get("value", "")))}</div>'
                        f'<div class="l">{escape(str(st.get("label", "")))}</div></div>'
                    )
                out.append("</div>")
            elif layout == "chart":
                out.append(f'<h2>{escape(s.get("title", ""))}</h2>')
                if s.get("subtitle"):
                    out.append(
                        f'<div style="font-size:13pt;color:#{muted};margin-top:-0.2in;'
                        f'margin-bottom:0.2in;">{escape(s["subtitle"])}</div>'
                    )
                chart_url = self._quickchart_url(s, p)
                out.append(f'<div class="chart-wrap"><img src="{escape(chart_url)}"/></div>')
                cd = s.get("chart_data") or {}
                cap_parts = []
                if cd.get("unit"):
                    cap_parts.append(f"Units: {escape(str(cd['unit']))}")
                if cd.get("source"):
                    cap_parts.append(f"Source: {escape(str(cd['source']))}")
                if cap_parts:
                    out.append(
                        f'<div class="chart-caption">{"  &bull;  ".join(cap_parts)}</div>'
                    )
            elif layout == "closing":
                out.append('<div class="closing">')
                out.append(f'<h1 class="title">{escape(s.get("title", "Thank you"))}</h1>')
                if s.get("subtitle"):
                    out.append(f'<div class="subtitle">{escape(s["subtitle"])}</div>')
                if s.get("cta"):
                    out.append(f'<div class="cta">{escape(s["cta"])}</div>')
                out.append("</div>")
            elif layout == "bigstat":
                # Phase 6AA — degrades to a stats-shaped HTML block so
                # any downstream HTML→PDF renderer can lay it out.
                out.append('<div class="bigstat">')
                if s.get("title"):
                    out.append(f'<div class="eyebrow">{escape(str(s["title"]))}</div>')
                out.append(
                    f'<div class="hero-value">{escape(str(s.get("value") or "—"))}</div>'
                )
                if s.get("label"):
                    out.append(f'<div class="hero-label">{escape(str(s["label"]))}</div>')
                if s.get("subtitle"):
                    out.append(f'<div class="subtitle">{escape(str(s["subtitle"]))}</div>')
                out.append("</div>")
            elif layout == "section_divider":
                # Phase 6AA — typography-only block; reuses .title CSS.
                out.append('<div class="section-divider">')
                if s.get("eyebrow"):
                    out.append(f'<div class="eyebrow">{escape(str(s["eyebrow"]))}</div>')
                out.append(f'<h1 class="title">{escape(s.get("title", ""))}</h1>')
                if s.get("subtitle"):
                    out.append(f'<div class="subtitle">{escape(str(s["subtitle"]))}</div>')
                out.append("</div>")
            elif layout == "timeline":
                # Phase 6AC — chronology rendered as a definition list so
                # any HTML→PDF renderer can lay it out without bespoke CSS.
                out.append('<div class="timeline">')
                out.append(f'<h1 class="title">{escape(s.get("title", ""))}</h1>')
                if s.get("subtitle"):
                    out.append(f'<div class="subtitle">{escape(str(s["subtitle"]))}</div>')
                events = s.get("events") or []
                if isinstance(events, list) and events:
                    out.append('<dl class="events">')
                    for e in events:
                        if not isinstance(e, dict):
                            continue
                        out.append(
                            f'<dt class="event-date">{escape(str(e.get("date") or ""))}</dt>'
                        )
                        out.append(
                            f'<dd class="event-label">{escape(str(e.get("label") or ""))}</dd>'
                        )
                    out.append("</dl>")
                out.append("</div>")
            elif layout == "comparison":
                # Phase 6AC — two-column comparison block with explicit
                # framing. Reuses .columns CSS so the layout downgrades
                # cleanly when no comparison-specific styles exist.
                out.append('<div class="comparison columns">')
                out.append(f'<h1 class="title">{escape(s.get("title", ""))}</h1>')
                if s.get("subtitle"):
                    out.append(f'<div class="subtitle">{escape(str(s["subtitle"]))}</div>')
                for side_key in ("left", "right"):
                    side = s.get(side_key) or {}
                    if not isinstance(side, dict):
                        continue
                    out.append(f'<div class="col col-{side_key}">')
                    if side.get("heading"):
                        out.append(
                            f'<h3 class="col-heading">{escape(str(side["heading"]))}</h3>'
                        )
                    if side.get("body"):
                        out.append(
                            f'<p class="col-body">{escape(str(side["body"]))}</p>'
                        )
                    out.append("</div>")
                out.append("</div>")
            out.append("</section>")
        out.append("</body></html>")
        return "".join(out)
