"""Extract slide structure from every PPTX in sample-decks/ for analysis.

Usage:
    python tools/analyze_manus_decks.py

Outputs one JSON per pptx beside it, plus a combined summary.json.
"""
from __future__ import annotations

import json
import sys
from pathlib import Path
from collections import Counter

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu
except ImportError:
    print("Installing python-pptx ...")
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu


ROOT = Path(__file__).resolve().parents[1] / "manus-reference" / "sample-decks"


def emu_to_px(v):
    if v is None:
        return None
    return round(Emu(v).inches * 96, 1)


def shape_info(shape) -> dict:
    info = {
        "name": shape.name,
        "type": str(shape.shape_type),
        "left": emu_to_px(shape.left),
        "top": emu_to_px(shape.top),
        "w": emu_to_px(shape.width),
        "h": emu_to_px(shape.height),
    }
    if shape.has_text_frame:
        runs = []
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                size = run.font.size.pt if run.font.size else None
                color = None
                try:
                    if run.font.color and run.font.color.rgb:
                        color = str(run.font.color.rgb)
                except Exception:
                    pass
                runs.append({
                    "text": run.text,
                    "bold": run.font.bold,
                    "italic": run.font.italic,
                    "size_pt": size,
                    "font": run.font.name,
                    "color": color,
                })
        info["text"] = shape.text_frame.text
        info["runs"] = runs
    if shape.has_table:
        tbl = shape.table
        info["table"] = {
            "rows": len(tbl.rows),
            "cols": len(tbl.columns),
            "cells": [[cell.text for cell in row.cells] for row in tbl.rows],
        }
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        info["picture"] = True
    return info


def classify_layout(shapes: list[dict]) -> str:
    """Heuristic layout label based on shape geometry + content."""
    text_shapes = [s for s in shapes if s.get("text", "").strip()]
    tables = [s for s in shapes if "table" in s]
    pictures = [s for s in shapes if s.get("picture")]

    n_text = len(text_shapes)
    n_table = len(tables)
    n_pic = len(pictures)

    # Title slide: 1-3 centered text shapes, no tables, no big body
    if n_text <= 3 and n_table == 0 and n_pic == 0:
        total_chars = sum(len(s.get("text", "")) for s in text_shapes)
        if total_chars < 200:
            return "title-cover"

    if n_table >= 1 and n_text >= 2 and n_pic == 0:
        # Side-by-side text + table
        if n_text >= 2:
            return "two-col-text-table"
        return "single-panel-with-table"

    if n_table == 0 and n_pic == 0:
        # All-text. Look at left positions to detect columns.
        lefts = sorted({round(s.get("left") or 0, -1) for s in text_shapes})
        if len(lefts) >= 2 and n_text >= 4:
            return "two-col-text-text"
        if n_text >= 4:
            return "stacked-panels"
        return "prose"

    if n_pic >= 3:
        return "icon-cards"
    if n_pic >= 1:
        return "image-focus"

    return "mixed"


def analyze_deck(pptx_path: Path) -> dict:
    prs = Presentation(str(pptx_path))
    slides = []
    layout_counts = Counter()
    color_counts = Counter()
    font_counts = Counter()
    size_counts = Counter()

    for idx, slide in enumerate(prs.slides, start=1):
        shapes = [shape_info(s) for s in slide.shapes]
        layout = classify_layout(shapes)
        layout_counts[layout] += 1

        for s in shapes:
            for r in s.get("runs", []):
                if r.get("color"):
                    color_counts[r["color"]] += 1
                if r.get("font"):
                    font_counts[r["font"]] += 1
                if r.get("size_pt"):
                    size_counts[round(r["size_pt"])] += 1

        # Plain-text dump of slide
        full_text = "\n".join(
            s.get("text", "").strip() for s in shapes if s.get("text", "").strip()
        )

        slides.append({
            "n": idx,
            "layout": layout,
            "n_shapes": len(shapes),
            "n_tables": sum(1 for s in shapes if "table" in s),
            "n_pictures": sum(1 for s in shapes if s.get("picture")),
            "word_count": len(full_text.split()),
            "char_count": len(full_text),
            "text": full_text,
            "shapes": shapes,
        })

    return {
        "file": str(pptx_path.relative_to(ROOT)),
        "slide_count": len(slides),
        "slide_dim_px": (
            emu_to_px(prs.slide_width),
            emu_to_px(prs.slide_height),
        ),
        "layout_distribution": dict(layout_counts),
        "top_colors": color_counts.most_common(10),
        "top_fonts": font_counts.most_common(10),
        "top_sizes_pt": size_counts.most_common(10),
        "avg_words_per_slide": round(
            sum(s["word_count"] for s in slides) / max(1, len(slides)), 1
        ),
        "slides": slides,
    }


def main():
    decks = list(ROOT.rglob("*.pptx"))
    print(f"Found {len(decks)} PPTX files")
    summary = []
    for path in decks:
        try:
            data = analyze_deck(path)
            out = path.with_suffix(".analysis.json")
            out.write_text(json.dumps(data, indent=2, ensure_ascii=False), encoding="utf-8")
            summary.append({
                "file": data["file"],
                "slides": data["slide_count"],
                "dim_px": data["slide_dim_px"],
                "layouts": data["layout_distribution"],
                "avg_words": data["avg_words_per_slide"],
                "top_colors": data["top_colors"][:5],
                "top_fonts": data["top_fonts"][:3],
            })
            print(f"  ✓ {path.name} → {data['slide_count']} slides, layouts={data['layout_distribution']}")
        except Exception as exc:
            print(f"  ✗ {path.name}: {exc}")

    (ROOT / "summary.json").write_text(
        json.dumps(summary, indent=2, ensure_ascii=False), encoding="utf-8"
    )
    print(f"\nSummary written to {ROOT / 'summary.json'}")


if __name__ == "__main__":
    main()
